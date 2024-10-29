import os
import gdown
from pptx import Presentation
import json
import logging
from src.ELearning.entity.config_entity import DataIngestionConfig

class DataIngestion:
    def __init__(self, config: DataIngestionConfig):
        self.config = config
        self.output_folder = os.path.join(self.config.root_dir)
        os.makedirs(self.output_folder, exist_ok=True)

    def download_pptx(self):
        logging.info("Downloading the presentation from Google Drive.")
        pptx_file_path = os.path.join(self.output_folder, self.config.output_filename)
        gdown.download(f'https://drive.google.com/uc?id={self.config.file_id}', pptx_file_path, quiet=False)
        logging.info(f"Downloaded PPTX file saved at {pptx_file_path}")
        return pptx_file_path

    def extract_text_and_images(self, pptx_file_path):
        prs = Presentation(pptx_file_path)
        extracted_data = []
        images_folder = os.path.join(self.output_folder, 'extracted_images')
        os.makedirs(images_folder, exist_ok=True)

        for slide_number, slide in enumerate(prs.slides):
            slide_data = {"slide_number": slide_number + 1, "text": "", "images": []}
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_data["text"] += shape.text + "\n"
                if shape.shape_type == 13:
                    image = shape.image
                    image_bytes = image.blob
                    image_filename = f"slide_{slide_number + 1}.{image.ext}"
                    image_path = os.path.join(images_folder, image_filename)
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    slide_data["images"].append(image_filename)

            extracted_data.append(slide_data)

        json_output_path = os.path.join(self.output_folder, "extracted_data.json")
        with open(json_output_path, "w") as json_file:
            json.dump(extracted_data, json_file, indent=4)
        logging.info(f"Extracted data saved to {json_output_path}")

    def initiate_data_ingestion(self):
        pptx_file_path = self.download_pptx()
        self.extract_text_and_images(pptx_file_path)
